"""
PowerPoint Presentation Generator
==================================

This script generates a professional PowerPoint presentation (.pptx)
for the Multi-Store Fashion Retail Sales Analysis.

Features:
- 16:9 widescreen format
- Modern professional design
- Embedded charts from reports/assets/
- Fully editable in PowerPoint, Google Slides, Keynote
- Consistent branding with blue color scheme

Author: PDF Reporting Specialist
Date: October 2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path
import pandas as pd


# Brand colors
DARK_BLUE = RGBColor(10, 64, 115)      # #0a4073
PRIMARY_BLUE = RGBColor(52, 152, 219)  # #3498db
LIGHT_BLUE = RGBColor(174, 214, 241)   # #aed6f1
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(44, 62, 80)       # #2c3e50
LIGHT_GRAY = RGBColor(236, 240, 241)   # #ecf0f1


def set_text_format(text_frame, font_name="Calibri", font_size=18,
                    bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    """
    Apply formatting to a text frame.

    Args:
        text_frame: TextFrame object
        font_name: Font family name
        font_size: Font size in points
        bold: Bold text flag
        color: RGB color object
        align: Text alignment
    """
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_title_slide(prs):
    """
    Add title slide.

    Args:
        prs: Presentation object
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add blue gradient background (simulated with shapes)
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height

    shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_BLUE

    # Title
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Multi-Store Fashion Retail\nSales Analysis"

    set_text_format(title_frame, font_size=44, bold=True,
                   color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    left = Inches(1)
    top = Inches(3.8)
    width = Inches(8)
    height = Inches(0.8)

    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "January 2024 Performance Analysis\nComplete 10-Store Dataset"

    set_text_format(subtitle_frame, font_size=24, bold=False,
                   color=WHITE, align=PP_ALIGN.CENTER)

    # Date
    left = Inches(1)
    top = Inches(5)
    width = Inches(8)
    height = Inches(0.5)

    date_box = slide.shapes.add_textbox(left, top, width, height)
    date_frame = date_box.text_frame
    date_frame.text = "October 20, 2025"

    set_text_format(date_frame, font_size=16, bold=False,
                   color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


def add_content_slide(prs, title, bullet_points):
    """
    Add a content slide with title and bullet points.

    Args:
        prs: Presentation object
        title: Slide title
        bullet_points: List of bullet point strings
    """
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    set_text_format(title_shape.text_frame, font_size=32, bold=True, color=DARK_BLUE)

    # Add bullet points
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = "Calibri"
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)


def add_chart_slide(prs, title, chart_path, caption=""):
    """
    Add a slide with a chart image.

    Args:
        prs: Presentation object
        title: Slide title
        chart_path: Path to chart image
        caption: Optional caption text
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title

    set_text_format(title_frame, font_size=32, bold=True,
                   color=DARK_BLUE, align=PP_ALIGN.LEFT)

    # Add chart image if it exists
    if Path(chart_path).exists():
        left = Inches(1)
        top = Inches(1.3)
        width = Inches(8)

        slide.shapes.add_picture(str(chart_path), left, top, width=width)

        # Add caption if provided
        if caption:
            left = Inches(1)
            top = Inches(5.2)
            width = Inches(8)
            height = Inches(0.4)

            caption_box = slide.shapes.add_textbox(left, top, width, height)
            caption_frame = caption_box.text_frame
            caption_frame.text = caption

            set_text_format(caption_frame, font_size=14, bold=False,
                           color=DARK_GRAY, align=PP_ALIGN.CENTER)


def add_metrics_slide(prs, sales_df):
    """
    Add a slide with key business metrics.

    Args:
        prs: Presentation object
        sales_df: Sales DataFrame
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Business Performance Overview"

    set_text_format(title_frame, font_size=32, bold=True,
                   color=DARK_BLUE, align=PP_ALIGN.LEFT)

    # Calculate metrics
    total_revenue = sales_df['sales_amount'].sum()
    total_transactions = len(sales_df)
    avg_transaction = sales_df['sales_amount'].mean()
    active_stores = sales_df['store_id'].nunique()

    # Metrics boxes
    metrics = [
        ("Total Revenue", f"¥{total_revenue/1000000:.1f}M"),
        ("Transactions", f"{total_transactions:,}"),
        ("Avg Transaction", f"¥{avg_transaction:,.0f}"),
        ("Active Stores", f"{active_stores} of 10")
    ]

    # Create 2x2 grid of metric boxes
    box_width = Inches(4)
    box_height = Inches(1.8)

    positions = [
        (Inches(0.8), Inches(1.5)),   # Top left
        (Inches(5.2), Inches(1.5)),   # Top right
        (Inches(0.8), Inches(3.5)),   # Bottom left
        (Inches(5.2), Inches(3.5))    # Bottom right
    ]

    for (left, top), (label, value) in zip(positions, metrics):
        # Background shape
        shape = slide.shapes.add_shape(1, left, top, box_width, box_height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_BLUE

        # Value (large number)
        value_box = slide.shapes.add_textbox(left, top + Inches(0.3), box_width, Inches(0.8))
        value_frame = value_box.text_frame
        value_frame.text = value
        set_text_format(value_frame, font_size=40, bold=True,
                       color=WHITE, align=PP_ALIGN.CENTER)

        # Label (description)
        label_box = slide.shapes.add_textbox(left, top + Inches(1.1), box_width, Inches(0.5))
        label_frame = label_box.text_frame
        label_frame.text = label
        set_text_format(label_frame, font_size=18, bold=False,
                       color=WHITE, align=PP_ALIGN.CENTER)


def add_recommendations_slide(prs):
    """
    Add strategic recommendations slide.

    Args:
        prs: Presentation object
    """
    recommendations = [
        "Replicate Fukuoka's strong 3rd-place performance in other regions",
        "Expand footwear inventory by 25-30% in Q2 (highest revenue category)",
        "Launch weekend activation initiatives to boost Saturday/Sunday sales",
        "Implement best practice exchange program between top and growth stores",
        "Balance Kanto concentration risk (39% share) with regional diversification"
    ]

    add_content_slide(prs, "Strategic Recommendations for Q2 2024", recommendations)


def add_closing_slide(prs):
    """
    Add closing slide.

    Args:
        prs: Presentation object
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height

    shape = slide.shapes.add_shape(1, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Thank you message
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1)

    message_box = slide.shapes.add_textbox(left, top, width, height)
    message_frame = message_box.text_frame
    message_frame.text = "Thank You"

    set_text_format(message_frame, font_size=48, bold=True,
                   color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    left = Inches(1)
    top = Inches(3.2)
    width = Inches(8)
    height = Inches(0.8)

    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions?\nReady for Q2 2024 Implementation"

    set_text_format(subtitle_frame, font_size=20, bold=False,
                   color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


def generate_powerpoint():
    """
    Main function to generate PowerPoint presentation.
    """
    print("=" * 80)
    print("POWERPOINT PRESENTATION GENERATOR")
    print("=" * 80)
    print()

    # Define paths
    project_root = Path(__file__).parent.parent.parent
    assets_dir = project_root / 'reports' / 'assets'
    output_file = project_root / 'reports' / 'executive_slides.pptx'
    sales_file = project_root / 'data' / 'processed' / 'sales_clean.csv'

    print(f"Assets directory: {assets_dir}")
    print(f"Output file: {output_file}")
    print()

    # Load sales data
    print("Loading sales data...")
    sales_df = pd.read_csv(sales_file)
    print(f"✓ Loaded {len(sales_df):,} transactions")
    print()

    # Create presentation
    print("Creating PowerPoint presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)    # 16:9 widescreen
    prs.slide_height = Inches(5.625)
    print("✓ Presentation initialized (16:9 format)")

    # Slide 1: Title
    print("\nGenerating slides:")
    print("  [1/13] Title slide")
    add_title_slide(prs)

    # Slide 2: Executive Summary
    print("  [2/13] Executive summary")
    exec_summary = [
        "Total revenue of ¥44.0M from 1,155 transactions across 10 stores",
        "Kanto region dominates with 39.2% share (4 stores including new Ikebukuro)",
        "Fukuoka debuts strong at #3 nationally, validating Kyushu expansion",
        "Footwear leads categories at 29.1% share (¥12.8M)",
        "Top 5 stores each contribute 10-12% revenue with balanced performance"
    ]
    add_content_slide(prs, "Executive Summary", exec_summary)

    # Slide 3: Business Overview
    print("  [3/13] Business metrics")
    add_metrics_slide(prs, sales_df)

    # Slide 4: Store Performance
    print("  [4/13] Store performance chart")
    chart_path = assets_dir / 'revenue_by_store.png'
    add_chart_slide(prs, "Store Performance Comparison", chart_path,
                   "Osaka leads with ¥5.2M, followed by Sendai and Fukuoka")

    # Slide 5: Regional Analysis
    print("  [5/13] Regional analysis chart")
    chart_path = assets_dir / 'revenue_by_region.png'
    add_chart_slide(prs, "Regional Revenue Distribution", chart_path,
                   "Kanto region generates ¥17.3M (39.2% share)")

    # Slide 6: Category Performance
    print("  [6/13] Category performance chart")
    chart_path = assets_dir / 'revenue_by_category.png'
    add_chart_slide(prs, "Product Category Performance", chart_path,
                   "Footwear and Accessories together account for 56.9% of revenue")

    # Slide 7: Weekday Patterns
    print("  [7/13] Day of week chart")
    chart_path = assets_dir / 'revenue_by_day_of_week.png'
    add_chart_slide(prs, "Sales Patterns by Day of Week", chart_path,
                   "Weekdays outperform weekends - opportunity for weekend activation")

    # Slide 8: Weekend vs Weekday
    print("  [8/13] Weekend vs weekday chart")
    chart_path = assets_dir / 'weekend_vs_weekday.png'
    add_chart_slide(prs, "Weekend vs Weekday Performance", chart_path,
                   "Weekdays generate 29% higher daily revenue")

    # Slide 9: Category Mix
    print("  [9/13] Category mix by store")
    chart_path = assets_dir / 'category_mix_by_store.png'
    add_chart_slide(prs, "Category Mix by Store", chart_path,
                   "All stores show balanced category distribution")

    # Slide 10: Top vs Bottom Performers
    print(" [10/13] Top vs bottom performers")
    chart_path = assets_dir / 'top_bottom_stores.png'
    add_chart_slide(prs, "Performance Gap Analysis", chart_path,
                   "42% gap between top and bottom performers presents opportunity")

    # Slide 11: Key Findings
    print(" [11/13] Key findings")
    findings = [
        "7-region national coverage achieved with 10-store network",
        "¥528M annual revenue run-rate (¥44M monthly)",
        "Consistent footwear strength across all stores indicates solid positioning",
        "Kyushu market entry successful - Fukuoka ranks #3 nationally",
        "Weekend activation opportunity identified"
    ]
    add_content_slide(prs, "Key Findings", findings)

    # Slide 12: Recommendations
    print(" [12/13] Recommendations")
    add_recommendations_slide(prs)

    # Slide 13: Closing
    print(" [13/13] Closing slide")
    add_closing_slide(prs)

    # Save presentation
    print()
    print("Saving PowerPoint presentation...")
    prs.save(str(output_file))

    file_size = output_file.stat().st_size / 1024
    print(f"✓ PowerPoint created: {output_file}")
    print(f"  File size: {file_size:.2f} KB")
    print()

    print("=" * 80)
    print("✅ POWERPOINT GENERATION COMPLETE")
    print("=" * 80)
    print()
    print(f"Presentation ready: {output_file}")
    print()
    print("You can now:")
    print("  - Open in Microsoft PowerPoint")
    print("  - Edit slides, rearrange, add animations")
    print("  - Present to executives")
    print("  - Open in Google Slides or Apple Keynote")
    print("  - Export to PDF from PowerPoint")
    print()

    return True


def main():
    """
    Main execution function.
    """
    try:
        success = generate_powerpoint()
        return success
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return False


if __name__ == "__main__":
    import sys
    success = main()
    sys.exit(0 if success else 1)
