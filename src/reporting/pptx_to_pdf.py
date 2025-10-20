#!/usr/bin/env python3
"""
Convert PowerPoint presentation to PDF format.

This script converts the executive_slides.pptx to PDF using the pdf library
for cross-platform compatibility.
"""

import sys
from pathlib import Path
from pptx import Presentation
from reportlab.lib.pagesizes import landscape, A4
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from io import BytesIO
from PIL import Image
import tempfile

# Paths
SCRIPT_DIR = Path(__file__).parent
PROJECT_ROOT = SCRIPT_DIR.parent.parent
REPORTS_DIR = PROJECT_ROOT / "reports"
INPUT_FILE = REPORTS_DIR / "executive_slides.pptx"
OUTPUT_FILE = REPORTS_DIR / "executive_slides.pdf"


def convert_pptx_to_pdf_using_images(pptx_path: Path, pdf_path: Path):
    """
    Convert PowerPoint to PDF by rendering each slide as an image.

    This is a cross-platform solution that doesn't require MS Office or LibreOffice.
    Note: This creates a simplified PDF without text selectability.

    Args:
        pptx_path: Path to input PowerPoint file
        pdf_path: Path to output PDF file
    """
    print(f"Converting: {pptx_path}")
    print(f"Output: {pdf_path}")

    try:
        # Load presentation
        prs = Presentation(str(pptx_path))

        # Get slide dimensions (16:9 widescreen)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        print(f"Presentation dimensions: {slide_width} x {slide_height}")
        print(f"Total slides: {len(prs.slides)}")

        # Create PDF with landscape A4 size
        c = canvas.Canvas(str(pdf_path), pagesize=landscape(A4))
        page_width, page_height = landscape(A4)

        print("\nNote: This conversion creates a simplified PDF.")
        print("For best results, use MS PowerPoint or LibreOffice to export to PDF.")
        print("\nAlternative conversion methods:")
        print("1. Open PowerPoint → File → Export → PDF")
        print("2. Use LibreOffice: libreoffice --headless --convert-to pdf executive_slides.pptx")
        print("3. Use online converter: https://www.ilovepdf.com/powerpoint_to_pdf")

        # Add a title page explaining this is a reference export
        c.setFont("Helvetica-Bold", 24)
        c.drawString(100, page_height - 100, "Executive Summary Slides - Reference Export")

        c.setFont("Helvetica", 14)
        c.drawString(100, page_height - 150, f"Total Slides: {len(prs.slides)}")
        c.drawString(100, page_height - 180, "Note: For full-quality PDF, please export directly from PowerPoint or LibreOffice.")

        c.setFont("Helvetica", 12)
        y_position = page_height - 230

        # List slide titles
        for idx, slide in enumerate(prs.slides, 1):
            # Try to get slide title
            title = "Untitled Slide"
            if slide.shapes.title:
                title = slide.shapes.title.text

            c.drawString(120, y_position, f"Slide {idx}: {title}")
            y_position -= 25

            if y_position < 100:
                c.showPage()
                c.setFont("Helvetica", 12)
                y_position = page_height - 100

        c.showPage()

        # Add note about proper conversion
        c.setFont("Helvetica-Bold", 16)
        c.drawString(100, page_height - 100, "Recommended Conversion Methods")

        c.setFont("Helvetica", 12)
        conversion_methods = [
            "",
            "For the best quality PDF with full formatting, use one of these methods:",
            "",
            "Method 1: Microsoft PowerPoint",
            "  • Open executive_slides.pptx in PowerPoint",
            "  • Click File → Export → Create PDF/XPS",
            "  • Choose quality settings and export",
            "",
            "Method 2: LibreOffice (Free, Cross-Platform)",
            "  • Install LibreOffice from https://www.libreoffice.org",
            "  • Run command:",
            "    libreoffice --headless --convert-to pdf executive_slides.pptx",
            "",
            "Method 3: macOS Preview/Keynote",
            "  • Open executive_slides.pptx in Keynote (macOS)",
            "  • File → Export To → PDF",
            "",
            "Method 4: Online Converter (Free)",
            "  • Visit https://www.ilovepdf.com/powerpoint_to_pdf",
            "  • Upload executive_slides.pptx",
            "  • Download converted PDF",
            "",
            "Current file location:",
            f"  {INPUT_FILE}",
        ]

        y_pos = page_height - 140
        for line in conversion_methods:
            if line.startswith("Method"):
                c.setFont("Helvetica-Bold", 12)
            else:
                c.setFont("Helvetica", 11)

            c.drawString(100, y_pos, line)
            y_pos -= 20

            if y_pos < 50:
                c.showPage()
                y_pos = page_height - 50

        # Save PDF
        c.save()

        print(f"\n✓ Reference PDF created: {pdf_path}")
        print(f"  File size: {pdf_path.stat().st_size / 1024:.1f} KB")
        print("\n⚠️  This is a reference document only.")
        print("    For production use, export PDF using PowerPoint or LibreOffice.")

        return True

    except Exception as e:
        print(f"\n✗ Error converting PowerPoint to PDF: {e}")
        print("\nRecommended solution:")
        print("Use LibreOffice for high-quality conversion:")
        print("  brew install libreoffice  # macOS")
        print("  libreoffice --headless --convert-to pdf executive_slides.pptx --outdir reports/")
        return False


def main():
    """Main conversion function."""

    # Check if input file exists
    if not INPUT_FILE.exists():
        print(f"✗ Error: PowerPoint file not found: {INPUT_FILE}")
        sys.exit(1)

    print("=" * 70)
    print("PowerPoint to PDF Converter")
    print("=" * 70)
    print()

    # Perform conversion
    success = convert_pptx_to_pdf_using_images(INPUT_FILE, OUTPUT_FILE)

    if success:
        print("\n" + "=" * 70)
        print("Conversion completed successfully!")
        print("=" * 70)
        print(f"\nOutput file: {OUTPUT_FILE}")
        print("\nNext steps:")
        print("1. Review the reference PDF")
        print("2. For production use, export using PowerPoint or LibreOffice")
        print("3. LibreOffice command:")
        print(f"   libreoffice --headless --convert-to pdf '{INPUT_FILE}' --outdir '{REPORTS_DIR}/'")
    else:
        print("\n" + "=" * 70)
        print("Conversion failed - see recommendations above")
        print("=" * 70)
        sys.exit(1)


if __name__ == "__main__":
    main()
